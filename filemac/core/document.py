"""Handler for dcoument conversion operations requested by the cli entry"""

import os
import re
import sqlite3
import subprocess
import sys

import pandas as pd
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pdf2docx import parse
from pdf2image import convert_from_path
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import Paragraph, SimpleDocTemplate
from rich.progress import Progress
from tqdm import tqdm
from ..utils.simple import logger
from ..utils.colors import fg, bg, rs

RESET = rs

DEFAULT_SEPARATOR = "\n"

_ext_word = ["doc", "docx"]
_ext_ppt_ = ["ppt", "pptx"]
_ext_xls = ["xls", "xlsx"]

PYGAME_DETECT_AVX2 = 1


class DocConverter:
    """Implementats all document conversion methods"""

    def __init__(self, input_file):
        self.input_file = input_file

    def preprocess(self):
        """Check input object whether it`s a file or a directory if a file append
        the file to a set and return it otherwise append directory full path
        content to the set and return the set file. The returned set will be
        evaluated in the next step as required on the basis of requested operation
        For every requested operation, the output file if any is automatically
        generated on the basis of the input filename and saved in the same
        directory as the input file.
        Exit if the folder is empty
        """

        try:
            files_to_process = []

            if os.path.isfile(self.input_file):
                files_to_process.append(self.input_file)
            elif os.path.isdir(self.input_file):
                if os.listdir(self.input_file) is None:
                    print("Cannot work with empty folder")
                    sys.exit(1)
                for file in os.listdir(self.input_file):
                    file_path = os.path.join(self.input_file, file)
                    if os.path.isfile(file_path):
                        files_to_process.append(file_path)

            return files_to_process
        except Exception as e:
            print(e)

    def word_to_pdf(self):
        """Convert word file to pdf document (docx)
        ->Check if running on Linux
        ->Use subprocess to run the dpkg and grep commands"""
        word_list = self.preprocess()

        word_list = [
            item for item in word_list if item.split(".")[-1].lower() in ("doc", "docx")
        ]
        for word_file in word_list:
            pdf_file_dir = os.path.dirname(word_file)
            pdf_file = os.path.splitext(word_file)[0] + ".pdf"

            try:
                if os.name == "posix":  # Check if running on Linux
                    print(
                        f"{fg.BLUE}Converting: {RESET}{word_file} {fg.BLUE}to {RESET}{pdf_file}"
                    )
                    # Use subprocess to run the dpkg and grep commands
                    result = subprocess.run(
                        ["dpkg", "-l", "libreoffice"], stdout=subprocess.PIPE, text=True
                    )
                    if result.returncode != 0:
                        logger.exception(f"{fg.RED}Libreoffice not found !{RESET}")
                        print(
                            f"{fg.CYAN}Initiating critical redundacy measure !{RESET}"
                        )
                        self.word2pdf_extra(word_file)
                    subprocess.run(
                        [
                            "soffice",
                            "--convert-to",
                            "pdf",
                            word_file,
                            "--outdir",
                            pdf_file_dir,
                        ]
                    )

                    print(
                        f"{fg.BMAGENTA} Successfully converted {word_file} to {pdf_file}{RESET}"
                    )
                    return pdf_file

                elif os.name == "nt":
                    self.word2pdf_extra(word_file)
                    return pdf_file

            except Exception as e:
                print(f"Error converting {word_file} to {pdf_file}: {e}")

    @staticmethod
    def word2pdf_extra(obj, outf=None):
        """For window users since it requires Microsoft word to be installed"""
        for file in obj:
            file = os.path.abspath(file)
            if file.split(".")[-1] not in ("doc", "docx"):
                logger.error(f"{fg.RED}File is not a word file{RESET}")
                sys.exit(1)
            pdf_file = os.path.splitext(file)[0] + ".pdf" if outf is None else outf
            try:
                if not os.path.isfile(file):
                    print(f"The file {obj} does not exist or is not a valid file.")
                    sys.exit("Exit!")
                logger.info(
                    f"{fg.BLUE}Converting: {RESET}{file} {fg.BLUE}to {RESET}{pdf_file}"
                )
                from docx2pdf import convert

                convert(file, pdf_file)
                print(f"{fg.GREEN}Conversion ✅{RESET}")
                sys.exit(0)
            except ImportError:
                logger.warning(
                    f"{fg.RED}docx2pdf Not found. {fg.CYAN}Run pip install docx2pdf{RESET}"
                )
            except Exception as e:
                raise
                logger.error(e)

    def pdf_to_word(self):
        """Convert pdf file to word document (docx)"""
        pdf_list = self.preprocess()
        pdf_list = [item for item in pdf_list if item.lower().endswith("pdf")]
        for pdf_file in pdf_list:
            word_file = (
                pdf_file[:-3] + "docx" if pdf_file.lower().endswith("pdf") else None
            )

            try:
                command = [
                    "lowriter",
                    "--headless",
                    '--infilter="writer_pdf_import"',
                    '--convert-todoc:"MS Word 97"',
                    pdf_file,
                ]
                print(f"{fg.BYELLOW}Parse the pdf document..{RESET}")
                parse(pdf_file, word_file, start=0, end=None)

                logger.info(
                    f"{fg.MAGENTA}New file is {fg.CYAN}{word_file}{RESET}"
                )
                logger.info(f"{fg.BGREEN}Success👨‍💻✅{RESET}")
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                logger.info(
                    f"{bg.RED}All conversion attempts have failed: {e}{RESET}"
                )

    def txt_to_pdf(self):
        """Convert text file(s) to pdf document (docx)
        ->Read the contents of the input .txt file
        ->Initialize the PDF document
        ->Create a story to hold the elements of the PDF
        ->Iterate through each line in the input .txt file and add it to the PDF
        ->Build and write the PDF document"""
        txt_list = self.preprocess()
        _list_ = [item for item in txt_list if item.lower().endswith("txt")]
        for _file_ in _list_:
            _pdf_ = _file_[:-3] + "pdf" if _file_.lower().endswith("txt") else None
            # Read the contents of the input .txt file
            with open(_file_, "r", encoding="utf-8") as file:
                text_contents = file.readlines()

            # Initialize the PDF document
            logger.info(f"{fg.BYELLOW}Initialize the PDF document{RESET}")
            doc = SimpleDocTemplate(_pdf_, pagesize=letter)

            # Create a story to hold the elements of the PDF
            logger.info(
                f"{fg.BYELLOW}Create a story to hold the elements of the PDF{RESET}"
            )
            story = []

            # Iterate through each line in the input .txt file and add it to the PDF
            logger.info(
                f"{fg.BYELLOW}Iterate through each line in the input .txt file and add it to the PDF{RESET}"
            )
            _line_count_ = 0
            try:
                for line in text_contents:
                    _line_count_ += 1
                    logger.info(
                        f"Lines {fg.BBLUE}{_line_count_}{RESET}/{len(text_contents)}"
                    )
                    story.append(Paragraph(line.strip(), style="normalText"))

            except KeyboardInterrupt:
                print("\nQuit❕⌨️")
                sys.exit(1)
            except Exception as e:
                logger.error(e)
                pass
            # Build and write the PDF document
            logger.info(f"{fg.BYELLOW}Build and write the PDF document{RESET}")
            doc.build(story)
            logger.info(f"{fg.MAGENTA}New file is {fg.CYAN}{_pdf_}{RESET}")
            print(f"\n{fg.BGREEN}Success👨‍💻✅{RESET}")

    def word_to_pptx(self):
        """Convert word file(s) to pptx document (pptx/ppt)
        -> Load the Word document
        ->Create a new PowerPoint presentation
        ->Iterate through each paragraph in the Word document
        ->Create a new slide in the PowerPoint presentation
        ->Add the paragraph text to the slide
        """
        word_list = self.preprocess()
        word_list = [
            item for item in word_list if item.split(".")[-1].lower() in ("doc", "docx")
        ]

        for word_file in word_list:
            if word_list is None:
                print("Please provide appropriate file type")
                sys.exit(1)
            ext = os.path.splitext(word_file)[-1][1:]

            pptx_file = (
                (os.path.splitext(word_file)[0] + ".pptx")
                if ext in list(_ext_word)
                else None
            )

            try:
                # Load the Word document
                print(f"{fg.BYELLOW}Load the Word document..{RESET}")
                doc = Document(word_file)

                # Create a new PowerPoint presentation
                print(f"{fg.BYELLOW}Create a new PowerPoint presentation..{RESET}")
                prs = Presentation()

                # Iterate through each paragraph in the Word document
                print(
                    f"{fg.BGREEN}Populating pptx slides with {fg.BYELLOW}{len(doc.paragraphs)}{fg.BGREEN} entries..{RESET}"
                )
                count = 0
                for paragraph in doc.paragraphs:
                    count += 1
                    perc = (count / len(doc.paragraphs)) * 100
                    print(
                        f"{fg.BMAGENTA}Progress:: {fg.BCYAN}{perc:.2f}%{RESET}",
                        end="\r",
                    )
                    # Create a new slide in the PowerPoint presentation
                    slide = prs.slides.add_slide(prs.slide_layouts[1])

                    # Add the paragraph text to the slide
                    slide.shapes.title.text = paragraph.text

                # Save the PowerPoint presentation
                prs.save(pptx_file)
                logger.info(
                    f"{fg.MAGENTA}New file is {fg.CYAN}{pptx_file}{RESET}"
                )
                print(f"\n{fg.BGREEN}Success👨‍💻✅{RESET}")
            except KeyboardInterrupt:
                print("\nQuit❕⌨️")
                sys.exit(1)
            except Exception as e:
                logger.error(e)

    def word_to_txt(self):
        """Convert word file to txt file"""
        word_list = self.preprocess()
        word_list = [
            item for item in word_list if item.split(".")[-1].lower() in ("dox", "docx")
        ]

        for file_path in word_list:
            ext = os.path.splitext(file_path)[-1][1:]
            txt_file = (
                (os.path.splitext(file_path)[0] + ".txt")
                if ext in list(_ext_word)
                else "output.txt"
            )

            try:
                logger.info(f"{fg.BLUE}Create Doument Tablet{RESET}")
                doc = Document(file_path)

                with open(txt_file, "w", encoding="utf-8") as f:
                    Par = 0
                    for paragraph in doc.paragraphs:
                        f.write(paragraph.text + "\n")
                        Par += 1

                        print(
                            f"Par:{fg.BLUE}{Par}/{len(doc.paragraphs)}{RESET}",
                            end="\r",
                        )
                    logger.info(
                        f"{fg.MAGENTA}Conversion of file to txt success{RESET}"
                    )

                logger.info(f"File: {fg.GREEN}{txt_file}{RESET}")
                return txt_file
            except KeyboardInterrupt:
                print("\nQuit❕⌨️")
                sys.exit()
            except Exception as e:
                logger.error(f"{fg.RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(
                        f"Couldn't convert {file_path} to {txt_file}:REASON->{e}"
                    )

    def pdf_to_txt(self):
        """Convert pdf file to text file"""

        pdf_list = self.preprocess()
        pdf_list = [item for item in pdf_list if item.lower().endswith("pdf")]
        for file_path in pdf_list:
            txt_file = file_path[:-3] + "txt"
            try:
                print(f"{fg.BYELLOW}Open and read the pdf document..{RESET}")
                with open(file_path, "rb") as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    _pg_ = 0
                    print(f"{fg.YELLOW}Convert pages..{RESET}")
                    for page_num in range(len(pdf_reader.pages)):
                        _pg_ += 1
                        logger.info(
                            f"Page {fg.BBLUE}{_pg_}{RESET}/{len(pdf_reader.pages)}"
                        )
                        page = pdf_reader.pages[page_num]
                        text += page.extract_text()
                with open(txt_file, "w", encoding="utf-8") as f:
                    f.write(text)
                logger.info(f"{fg.MAGENTA}New file is {fg.CYAN}{txt_file}{RESET}")
                logger.info(f"{fg.BGREEN}Success👨‍💻✅{RESET}")
            except Exception as e:
                logger.error(f"{fg.RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(f"Error converting {file_path} to {txt_file}: {e}\n")

    def pptx_to_txt(self, dest=None):
        """Convert ppt file to tetx document"""
        ppt_list = self.preprocess()
        ppt_list = [
            item for item in ppt_list if item.split(".")[-1].lower() in ("ppt", "pptx")
        ]
        try:
            for file_path in ppt_list:
                ext = os.path.splitext(file_path)[-1][1:]

                txt_file = (os.path.splitext(file_path)[0]) + ".txt"

                file_path = os.path.abspath(file_path)

                if ext == "ppt":
                    file_path = self.convert_ppt_to_pptx(
                        file_path
                    )  # First convert the ppt to pptx

                presentation = Presentation(file_path)

                logger.info(
                    f"Slide count ={fg.BMAGENTA} {len(presentation.slides)}{RESET}"
                )

                _slide_count_ = 0

                with Progress() as progress:
                    task = progress.add_task(
                        "[magenta]Preparing..", total=len(presentation.slides)
                    )

                    for slide in presentation.slides:
                        _slide_count_ += 1
                        # progress.console.print(F"Slide {_slide_count_}/{len(presentation.slides)}", end='\n')

                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame

                                for paragraph in text_frame.paragraphs:
                                    # Create a paragraph in the Word document if it contains text
                                    # Ensure text exists
                                    if any(run.text.strip() for run in paragraph.runs):
                                        for run in paragraph.runs:
                                            text = run.text.strip()
                                            if text and text != " ":
                                                with open(txt_file, "a") as fl:
                                                    fl.write(text)
                                                    # return txt_file

                        progress.update(task, advance=1)

                    if dest == "text":
                        with open(txt_file, "r") as fl:
                            text_buffer = fl.read()
                        print(text_buffer)
                        return text_buffer

                logger.info(f"{fg.MAGENTA}New file is {fg.CYAN}{txt_file}{RESET}")
                logger.info(f"{fg.BGREEN}Success👨‍💻✅{RESET}")
        except Exception as e:
            logger.error(f"\n❌Oops! {bg.RED}{e}{RESET}")

    @staticmethod
    def convert_ppt_to_pptx(obj: os.PathLike):
        import platform

        try:
            if obj.endswith("ppt"):
                if platform.system() in ("Linux", "MacOS") or os.name == "posix":
                    subprocess.run(
                        ["soffice", "--headless", "--convert-to", "pptx", obj]
                    )
                    return os.path.splitext(obj)[0] + ".pptx"
                elif platform.system() in ("Windows") or os.name == "nt":
                    import win32com.client

                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    powerpoint.Visible = 1
                    ppt = powerpoint.Presentations.Open(obj)
                    pptx_file = os.path.splitext(obj)[0] + ".pptx"
                    ppt.SaveAs(pptx_file, 24)  # 24 is the format for pptx
                    ppt.Close()
                    powerpoint.Quit()
                    return pptx_file
            else:
                print(f"{fg.RED}Unable to identify the system{RESET}")
        except KeyboardInterrupt:
            print("\nQuit!")
            sys.exit(1)
        except Exception as e:
            logger.error(f"{fg.RED}{e}{RESET}")

    def ppt_to_word(self):
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        from docx.shared import Pt
        from docx.shared import RGBColor as docxRGBColor
        from pptx.dml.color import RGBColor as pptxRGBColor

        """Convert ppt file to word document\n
        ->Preserves bold formatting
        """
        ppt_list = self.preprocess()
        ppt_list = [
            item for item in ppt_list if item.split(".")[-1].lower() in ("ppt", "pptx")
        ]
        for file_path in ppt_list:
            ext = os.path.splitext(file_path)[-1][1:]
            word_file = (
                (os.path.splitext(file_path)[0] + ".docx")
                if ext in list(_ext_ppt_)
                else None
            )
            try:
                logger.info(f"{fg.BYELLOW}Create Doument Tablet{RESET}")
                file_path = os.path.abspath(file_path)
                if ext == "ppt":
                    file_path = self.convert_ppt_to_pptx(
                        file_path
                    )  # First convert the ppt to pptx
                presentation = Presentation(file_path)
                document = Document()
                logger.info(
                    f"Slide count ={fg.BMAGENTA} {len(presentation.slides)}{RESET}"
                )
                _slide_count_ = 0
                with Progress() as progress:
                    task = progress.add_task(
                        "[magenta]Preparing..", total=len(presentation.slides)
                    )
                    for slide in presentation.slides:
                        _slide_count_ += 1
                        # progress.console.print(F"Slide {_slide_count_}/{len(presentation.slides)}", end='\n')
                        slide_text = ""
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame

                                for paragraph in text_frame.paragraphs:
                                    # Create a paragraph in the Word document if it contains text
                                    # Ensure text exists
                                    if any(run.text.strip() for run in paragraph.runs):
                                        # print("Has text")
                                        new_paragraph = document.add_paragraph()

                                        # Set general paragraph properties
                                        new_paragraph.alignment = (
                                            WD_PARAGRAPH_ALIGNMENT.JUSTIFY
                                        )  # Justify text
                                        new_paragraph.space_after = Pt(6)
                                        new_paragraph.space_before = Pt(6)
                                        new_paragraph.line_spacing = 1.15

                                        for run in paragraph.runs:
                                            if run.text.strip():
                                                slide_text += (
                                                    run.text
                                                )  # Only add non-empty text runs
                                                # print(run.text.strip(), end='\n')
                                                new_run = new_paragraph.add_run(
                                                    run.text
                                                )

                                                # Preserve bold, italic, underline, font name, and size
                                                new_run.bold = run.font.bold
                                                new_run.italic = run.font.italic
                                                new_run.underline = run.font.underline
                                                new_run.font.name = run.font.name
                                                new_run.font.size = run.font.size

                                                # Preserve font color
                                                try:
                                                    if (
                                                        run.font.color
                                                        and run.font.color.rgb
                                                    ):
                                                        pptx_color = run.font.color.rgb
                                                        # If the color is white (255, 255, 255), change it to black (0, 0, 0)
                                                        if pptx_color == pptxRGBColor(
                                                            255, 255, 255
                                                        ):
                                                            new_run.font.color.rgb = (
                                                                docxRGBColor(0, 0, 0)
                                                            )  # Black
                                                        else:
                                                            # Assign color properly to the Word run
                                                            new_run.font.color.rgb = (
                                                                docxRGBColor(
                                                                    pptx_color[0],
                                                                    pptx_color[1],
                                                                    pptx_color[2],
                                                                )
                                                            )
                                                except AttributeError:
                                                    pass

                        progress.update(task, advance=1)
                document.save(word_file)
                logger.info(
                    f"{fg.MAGENTA}New file is {fg.CYAN}{word_file}{RESET}"
                )
                logger.info(f"{fg.BGREEN}Success👨‍💻✅{RESET}")
                return word_file
            except Exception as e:
                logger.error(f"\n❌Oops! {bg.RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(f"\n❌Oops! {e}")

    def text_to_word(self):
        """Convert text file to word\n
        ->Read the text file\n
        ->Filter out non-XML characters\n
        ->Create a new Word document\n
        ->Add the filtered text content to the document"""
        flist = self.preprocess()
        flist = [item for item in flist if item.lower().endswith("txt")]
        for file_path in flist:
            if file_path.lower().endswith("txt"):
                word_file = file_path[:-3] + "docx"

            try:
                # Read the text file
                logger.info(f"{fg.BCYAN}Open and read the text file{RESET}")
                with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
                    text_content = file.read()

                # Filter out non-XML characters
                filtered_content = re.sub(
                    r"[^\x09\x0A\x0D\x20-\uD7FF\uE000-\uFFFD]+", "", text_content
                )

                # Create a new Word document
                logger.info(f"{fg.BYELLOW}Create Doument Tablet{RESET}")
                doc = Document()
                # Add the filtered text content to the document
                doc.add_paragraph(filtered_content)

                # Save the document as a Word file
                doc.save(word_file)
                logger.info(
                    f"{fg.MAGENTA}New file is {fg.BCYAN}{word_file}{RESET}"
                )
                logger.info(f"{fg.BGREEN}Success👨‍💻✅{RESET}")
            except FileExistsError as e:
                logger.error(f"{str(e)}📁")
            except Exception as e:
                logger.error(f"\n❌Oops something went awry {fg.RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(
                        f"\n❌Oops something went astray{fg.RED}{e}{RESET}"
                    )

    def convert_xls_to_word(self):
        """Convert xlsx file(s) to word file(s)\n
        ->Read the XLS file using pandas\n
        ->Create a new Word document\n
        ->Iterate over the rows of the dataframe and add them to the Word document"""
        xls_list = self.preprocess()

        xls_list = [
            item for item in xls_list if item.split(".")[-1].lower() in ("xls", "xlsx")
        ]

        print(f"{fg.BGREEN}Initializing conversion sequence{RESET}")

        for xls_file in xls_list:
            ext = os.path.splitext(xls_file)[-1][1:]
            word_file = (
                (os.path.splitext(xls_file)[0] + ".docx")
                if ext in list(_ext_xls)
                else None
            )
            try:
                """Read the XLS file using pandas"""

                df = pd.read_excel(xls_file)

                """Create a new Word document"""
                doc = Document()

                """Iterate over the rows of the dataframe and add them to the
                Word document"""
                logger.info(f"{fg.ICYAN}Converting {xls_file}..{RESET}")
                # time.sleep(2)
                total_rows = df.shape[0]
                for _, row in df.iterrows():
                    current_row = _ + 1
                    percentage = (current_row / total_rows) * 100
                    for value in row:
                        doc.add_paragraph(str(value))
                    print(
                        f"Row {fg.BYELLOW}{current_row}/{total_rows} {fg.BBLUE}{percentage:.1f}%{RESET}",
                        end="\r",
                    )
                    # print(f"\033[1;36m{row}{RESET}")

                # Save the Word document
                doc.save(word_file)
                print(f"{fg.BGREEN}Conversion successful!{RESET}", end="\n")
            except KeyboardInterrupt:
                print("\nQuit⌨️")
                sys.exit(1)
            except Exception as e:
                print(f"{bg.RED}Oops Conversion failed:❕{RESET}", str(e))

    def convert_xls_to_text(self):
        """Convert xlsx/xls file/files to text file format
        ->Read the XLS file using pandas
        ->Convert the dataframe to plain text
        ->Write the plain text to the output file"""
        xls_list = self.preprocess()

        xls_list = [
            item
            for item in xls_list
            if any(item.lower().endswith(ext) for ext in _ext_xls)
        ]
        print(f"{fg.BGREEN}Initializing conversion sequence{RESET}")
        for xls_file in tqdm(xls_list):
            ext = os.path.splitext(xls_file)[-1][1:]
            txt_file = (
                (os.path.splitext(xls_file)[0] + ".txt")
                if ext in list(_ext_xls)
                else None
            )
            try:
                # Read the XLS file using pandas
                logger.info(f"Converting {xls_file}..")
                df = pd.read_excel(xls_file)

                # Convert the dataframe to plain text
                text = df.to_string(index=False)
                chars = len(text)
                words = len(text.split())
                lines = len(text.splitlines())

                print(
                    f"Preparing to write: {fg.BYELLOW}{chars} \033[1;30m characters{fg.BYELLOW} {words}\033[1;30m words {fg.BYELLOW}{lines}\033[1;30m lines {RESET}",
                    end="\n",
                )
                # Write the plain text to the output file
                with open(txt_file, "w") as file:
                    file.write(text)

                print(f"{fg.BGREEN}Conversion successful!{RESET}", end="\n")
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                print("Oops Conversion failed:", str(e))

    def convert_xlsx_to_csv(self):
        """Convert xlsx/xls file to csv(comma seperated values) format
        ->Load the Excel file
        ->Save the DataFrame to CSV"""
        xls_list = self.preprocess()

        xls_list = [
            item for item in xls_list if item.split(".")[-1].lower() in ("xls", "xlsx")
        ]
        for xls_file in tqdm(xls_list):
            ext = os.path.splitext(xls_file)[-1][1:]
            csv_file = (
                (os.path.splitext(xls_file)[0] + ".csv")
                if ext in list(_ext_xls)
                else None
            )
            try:
                """Load the Excel file"""
                print(f"{fg.BGREEN}Initializing conversion sequence{RESET}")
                df = pd.read_excel(xls_file)
                logger.info(f"Converting {xls_file}..")
                total_rows = df.shape[0]
                print(f"Writing {fg.BYELLOW}{total_rows} rows {RESET}", end="\n")
                for i in range(101):
                    print(f"Progress: {i}%", end="\r")
                """Save the DataFrame to CSV"""
                df.to_csv(csv_file, index=False)
                print(f"{fg.BMAGENTA} Conversion successful{RESET}")

            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                print(e)

    def convert_csv_to_xlsx(self):
        csv_list = self.preprocess()
        csv_list = [item for item in csv_list if item.split(".")[-1].lower() in ("csv")]

        with Progress() as progress:
            task = progress.add_task("[cyan]Coverting", total=len(csv_list))
            for file in csv_list:
                file_name = file[:-3] + "xlsx"
                df = pd.read_csv(file)
                # excel engines ('openpyxl' or 'xlsxwriter')
                df.to_excel(file_name, engine="openpyxl", index=False)

                # Load the workbook and the sheet
                workbook = load_workbook(file_name)
                sheet = workbook.active

                # print("Adjust Columns")
                for column in sheet.columns:
                    max_length = 0
                    column_letter = column[0].column_letter

                    max_length = max(len(str(cell.value)) for cell in column)
                    """for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(cell.value)

                        except Exception:
                            pass
                    """

                    adjusted_width = max_length + 2
                    sheet.column_dimensions[column_letter].width = adjusted_width

                # Save the workbook
                workbook.save(file_name)
                progress.update(task, advance=1)

    def convert_xlsx_to_database(self):
        """Convert xlsx file(s) to sqlite
        ->Read the Excel file into a pandas DataFrame
        ->Create a connection to the SQLite database
        ->Insert the DataFrame into a new table in the database
        ->Close the database connection"""
        xlsx_list = self.preprocess()
        xlsx_list = [
            item for item in xlsx_list if item.split(".")[-1].lower() in ("xls", "xlsx")
        ]
        for xlsx_file in tqdm(xlsx_list):
            sqlfile = (
                (os.path.splitext(xlsx_file)[0] + ".sql")
                if (xlsx_file.split(".")[0]) in ("xls", "xlsx")
                else None
            )
            try:
                db_file = input(
                    f"{fg.BBLUE}Please enter desired sql filename: {RESET}"
                )
                table_name = input("Please enter desired table name: ")
                # res = ["db_file", "table_name"]
                if any(db_file) == "":
                    db_file = sqlfile
                    table_name = sqlfile[:-4]
                if not db_file.endswith(".sql"):
                    db_file = sqlfile
                column = 0
                for i in range(20):
                    column += 0
                # Read the Excel file into a pandas DataFrame
                print(f"Reading {xlsx_file}...")
                df = pd.read_excel(xlsx_file)
                print(f"{fg.BGREEN}Initializing conversion sequence{RESET}")
                print(f"{fg.BGREEN} Connected to sqlite3 database::{RESET}")
                # Create a connection to the SQLite database
                conn = sqlite3.connect(db_file)
                print(f"{fg.BYELLOW} Creating database table::{RESET}")
                # Insert the DataFrame into a new table in the database
                df.to_sql(table_name, column, conn, if_exists="replace", index=False)
                print(
                    f"Operation successful{RESET} file saved as \033[32{db_file}{RESET}"
                )
                # Close the database connection
                conn.close()
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                logger.error(f"{e}")

    def doc2image(self, outf="png"):
        """Create image objects from given files"""
        outf = "png" if outf not in ("png", "jpg") else outf
        path_list = self.preprocess()
        file_list = [
            item
            for item in path_list
            if item.split(".")[-1].lower() in ("pdf", "doc", "docx")
        ]
        imgs = []
        for file in file_list:
            if file.lower().endswith("pdf"):
                # Convert the PDF to a list of PIL image objects
                print(f"{fg.BLUE}Generate image objects ..{RESET}")
                images = convert_from_path(file)

                # Save each image to a file
                fname = file[:-4]
                print(f"{fg.YELLOW}Target images{fg.BLUE} {len(images)}{RESET}")

                with Progress() as progress:
                    task = progress.add_task(
                        "[magenta]Generating images ", total=len(images)
                    )
                    for i, image in enumerate(images):
                        # print(f"{Bfg.BLUE}{i}{RESET}", end="\r")
                        yd = f"{fname}_{i + 1}.{outf}"
                        image.save(yd)
                        imgs.append(yd)
                        progress.update(task, advance=1)
                # print(f"\n{fg.GREEN}Ok{RESET}")

        return imgs
