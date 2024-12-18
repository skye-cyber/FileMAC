#############################################################################
import json
import logging
import logging.handlers
import math
import os
import re
import shutil
import sqlite3
import subprocess
import sys
from threading import Lock, Thread
from typing import List, Union

import cv2
import pandas as pd
import pydub
import PyPDF2
import requests
from docx import Document
from gtts import gTTS
from moviepy.editor import VideoFileClip
from openpyxl import load_workbook
from pdf2docx import parse
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pydub import AudioSegment
from reportlab.lib.pagesizes import letter
from reportlab.platypus import Paragraph, SimpleDocTemplate
from rich.errors import MarkupError
from rich.progress import Progress
from tqdm import tqdm

# from pathlib import Path
from .colors import (BLUE, BWHITE, CYAN, DBLUE, DCYAN, DGREEN, DMAGENTA, DRED,
                     DYELLOW, FCYAN, FMAGENTA, GREEN, ICYAN, IGREEN, MAGENTA,
                     RED, RESET, YELLOW)
from .formats import (SUPPORTED_AUDIO_FORMATS, SUPPORTED_AUDIO_FORMATS_DIRECT,
                      SUPPORTED_IMAGE_FORMATS, SUPPORTED_VIDEO_FORMATS,
                      Video_codecs)
from .m4a_converter import _m4a_main_

# import pygame
# from aspose.words import Document as aspose_document
# from aspose.slides import Presentation as aspose_presentation
# from show_progress import progress_show
# from PIL import ImageDraw, ImageFont
###############################################################################

_ext_word = ["doc", "docx"]
_ext_ppt_ = ["ppt", "pptx"]
_ext_xls = ["xls", "xlsx"]

PYGAME_DETECT_AVX2 = 1
logging.basicConfig(level=logging.INFO, format='%(levelname)-8s %(message)s')
logger = logging.getLogger(__name__)


class MakeConversion:

    '''Implementation for all file conversions'''

    def __init__(self, input_file):
        self.input_file = input_file

    def preprocess(self):
        '''Check input object whether it`s a file or a directory if a file append
    the file to a set and return it otherwise append directory full path
    content to the set and return the set file. The returned set will be
    evaluated in the next step as required on the basis of requested operation
    For every requested operation, the output file if any is automatically
    generated on the basis of the input filename and saved in the same
    directory as the input file.
    Exit if the folder is empty
    '''

        try:
            files_to_process = []

            if os.path.isfile(self.input_file):
                files_to_process.append(self.input_file)
            elif os.path.isdir(self.input_file):
                if os.listdir(self.input_file) is None:
                    print("Cannot work with empty folder")
                    sys.exit(1)
                for file in os.listdir(self.input_file):
                    file_path = os.path.join(self.input_file, file)
                    if os.path.isfile(file_path):
                        files_to_process.append(file_path)

            return files_to_process
        except Exception as e:
            print(e)

    def word_to_pdf(self):
        ###############################################################################
        '''Convert word file to pdf document (docx)
        ->Check if running on Linux
        ->Use subprocess to run the dpkg and grep commands'''
        ###############################################################################
        word_list = self.preprocess()

        word_list = [
            item for item in word_list if item.split('.')[-1].lower() in ("doc", "docx")]
        for word_file in word_list:

            pdf_file_dir = os.path.dirname(word_file)
            pdf_file = os.path.splitext(word_file)[0] + '.pdf'

            try:
                if os.name == 'posix':  # Check if running on Linux
                    print(
                        f'{BLUE}Converting: {RESET}{word_file} {BLUE}to {RESET}{pdf_file}')
                    # Use subprocess to run the dpkg and grep commands
                    result = subprocess.run(
                        ['dpkg', '-l', 'libreoffice'], stdout=subprocess.PIPE, text=True)
                    if result.returncode != 0:
                        logger.exception(
                            f"{RED}Libreoffice not found !{RESET}")
                        print(
                            f"{CYAN}Initiating critical redundacy measure !{RESET}")
                        self.word2pdf_extra(word_file)
                    subprocess.run(
                        ['soffice', '--convert-to', 'pdf',
                         word_file, '--outdir', pdf_file_dir])

                    print(
                        f"{DMAGENTA} Successfully converted {word_file} to {pdf_file}{RESET}")
                    return pdf_file

                elif os.name == "nt":
                    self.word2pdf_extra(word_file)
                    return pdf_file

            except Exception as e:
                print(f"Error converting {word_file} to {pdf_file}: {e}")

    @staticmethod
    def word2pdf_extra(obj, outf=None):
        '''For window users since it requires Microsoft word to be installed'''
        for file in obj:
            file = os.path.abspath(file)
            if file.split('.')[-1] not in ('doc', 'docx'):
                logger.error(f"{RED}File is not a word file{RESET}")
                sys.exit(1)
            pdf_file = os.path.splitext(
                file)[0] + '.pdf' if outf is None else outf
            try:
                if not os.path.isfile(file):
                    print(
                        f"The file {obj} does not exist or is not a valid file.")
                    sys.exit('Exit!')
                logger.info(
                    f'{BLUE}Converting: {RESET}{file} {BLUE}to {RESET}{pdf_file}')
                from docx2pdf import convert
                convert(file, pdf_file)
                print(F"{GREEN}Conversion ✅{RESET}")
                sys.exit(0)
            except ImportError:
                logger.warning(f"{RED}docx2pdf Not found. {
                    CYAN}Run pip install docx2pdf{RESET}")
            except Exception as e:
                raise
                logger.error(e)

    def pdf_to_word(self):
        ###############################################################################
        """Convert pdf file to word document (docx)"""
        ###############################################################################
        pdf_list = self.preprocess()
        pdf_list = [item for item in pdf_list if item.lower().endswith("pdf")]
        for pdf_file in pdf_list:
            word_file = pdf_file[:-3] + \
                "docx" if pdf_file.lower().endswith("pdf") else None

            try:
                print(F"{DYELLOW}Parse the pdf document..{RESET}")
                parse(pdf_file, word_file, start=0, end=None)

                logger.info(f"{MAGENTA}New file is {CYAN}{word_file}{RESET}")
                logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                logger.info(f'{DRED}All conversion attempts have failed: \
{e}{RESET}')

    def txt_to_pdf(self):
        ###############################################################################
        """Convert text file(s) to pdf document (docx)
        ->Read the contents of the input .txt file
        ->Initialize the PDF document
        ->Create a story to hold the elements of the PDF
        ->Iterate through each line in the input .txt file and add it to the PDF
        ->Build and write the PDF document"""
        ###############################################################################
        txt_list = self.preprocess()
        _list_ = [item for item in txt_list if item.lower().endswith("txt")]
        for _file_ in _list_:
            _pdf_ = _file_[:-3] + \
                "pdf" if _file_.lower().endswith("txt") else None
            # Read the contents of the input .txt file
            with open(_file_, 'r', encoding='utf-8') as file:
                text_contents = file.readlines()

            # Initialize the PDF document
            logger.info(F"{DYELLOW}Initialize the PDF document{RESET}")
            doc = SimpleDocTemplate(_pdf_, pagesize=letter)

            # Create a story to hold the elements of the PDF
            logger.info(
                F"{DYELLOW}Create a story to hold the elements of the PDF{RESET}")
            story = []

            # Iterate through each line in the input .txt file and add it to the PDF
            logger.info(
                F"{DYELLOW}Iterate through each line in the input .txt file and add it to the PDF{RESET}")
            _line_count_ = 0
            try:
                for line in text_contents:
                    _line_count_ += 1
                    logger.info(
                        f"Lines {DBLUE}{_line_count_}{RESET}/{len(text_contents)}")
                    story.append(Paragraph(line.strip(), style="normalText"))

            except KeyboardInterrupt:
                print("\nQuit❕⌨️")
                sys.exit(1)
            except Exception as e:
                logger.error(e)
                pass
            # Build and write the PDF document
            logger.info(f"{DYELLOW}Build and write the PDF document{RESET}")
            doc.build(story)
            logger.info(f"{MAGENTA}New file is {CYAN}{_pdf_}{RESET}")
            print(f"\n{DGREEN}Success👨‍💻✅{RESET}")

    def word_to_pptx(self):
        ###############################################################################
        """Convert word file(s) to pptx document (pptx/ppt)
        -> Load the Word document
        ->Create a new PowerPoint presentation
        ->Iterate through each paragraph in the Word document
        ->Create a new slide in the PowerPoint presentation
        ->Add the paragraph text to the slide
        """
        ###############################################################################
        word_list = self.preprocess()
        word_list = [item for item in word_list if item.split(
            '.')[-1].lower() in ("doc", "docx")]

        for word_file in word_list:

            if word_list is None:
                print("Please provide appropriate file type")
                sys.exit(1)
            ext = os.path.splitext(word_file)[-1][1:]

            pptx_file = (os.path.splitext(word_file)[
                         0] + ".pptx") if ext in list(_ext_word) else None

            try:
                # Load the Word document
                print(F"{DYELLOW}Load the Word document..{RESET}")
                doc = Document(word_file)

                # Create a new PowerPoint presentation
                print(F"{DYELLOW}Create a new PowerPoint presentation..{RESET}")
                prs = Presentation()

                # Iterate through each paragraph in the Word document
                print(
                    f"{DGREEN}Populating pptx slides with {DYELLOW}{len(doc.paragraphs)}{DGREEN} entries..{RESET}")
                count = 0
                for paragraph in doc.paragraphs:
                    count += 1
                    perc = (count/len(doc.paragraphs))*100
                    print(
                        f"{DMAGENTA}Progress:: {DCYAN}{perc:.2f}%{RESET}", end="\r")
                    # Create a new slide in the PowerPoint presentation
                    slide = prs.slides.add_slide(prs.slide_layouts[1])

                    # Add the paragraph text to the slide
                    slide.shapes.title.text = paragraph.text

                # Save the PowerPoint presentation
                prs.save(pptx_file)
                logger.info(f"{MAGENTA}New file is {CYAN}{pptx_file}{RESET}")
                print(f"\n{DGREEN}Success👨‍💻✅{RESET}")
            except KeyboardInterrupt:
                print("\nQuit❕⌨️")
                sys.exit(1)
            except Exception as e:
                logger.error(e)

    def word_to_txt(self):
        ###############################################################################
        """Convert word file to txt file"""
        ###############################################################################
        word_list = self.preprocess()
        word_list = [item for item in word_list if item.split(
            '.')[-1].lower() in ("dox", "docx")]

        for file_path in word_list:
            ext = os.path.splitext(file_path)[-1][1:]
            txt_file = (os.path.splitext(file_path)[
                        0] + ".txt") if ext in list(_ext_word) else "output.txt"

            try:
                logger.info(f"{BLUE}Create Doument Tablet{RESET}")
                doc = Document(file_path)

                with open(txt_file, 'w', encoding='utf-8') as f:
                    Par = 0
                    for paragraph in doc.paragraphs:
                        f.write(paragraph.text + '\n')
                        Par += 1

                        print(
                            f"Par:{BLUE}{Par}/{len(doc.paragraphs)}{RESET}", end='\r')
                    logger.info(
                        f"{MAGENTA}Conversion of file to txt success{RESET}")

                logger.info(f"File: {GREEN}{txt_file}{RESET}")
                return txt_file
            except KeyboardInterrupt:
                print("\nQuit❕⌨️")
                sys.exit()
            except Exception as e:
                logger.error(f"{RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(f"Couldn't convert {file_path} to {txt_file}:\
REASON->{e}")

    def pdf_to_txt(self):

        ###############################################################################
        """Convert pdf file to text file"""
        ###############################################################################

        pdf_list = self.preprocess()
        pdf_list = [item for item in pdf_list if item.lower().endswith("pdf")]
        for file_path in pdf_list:
            txt_file = file_path[:-3] + "txt"
            try:
                print(F"{DYELLOW}Open and read the pdf document..{RESET}")
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ''
                    _pg_ = 0
                    print(F"{YELLOW}Convert pages..{RESET}")
                    for page_num in range(len(pdf_reader.pages)):
                        _pg_ += 1
                        logger.info(
                            f"Page {DBLUE}{_pg_}{RESET}/{len(pdf_reader.pages)}")
                        page = pdf_reader.pages[page_num]
                        text += page.extract_text()
                with open(txt_file, 'w', encoding='utf-8') as f:
                    f.write(text)
                logger.info(f"{MAGENTA}New file is {CYAN}{txt_file}{RESET}")
                logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            except Exception as e:
                logger.error(f"{RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(
                        f"Error converting {file_path} to {txt_file}: {e}\n")

    def pptx_to_txt(self, dest=None):
        ###############################################################################
        """Convert ppt file to tetx document"""
        ###############################################################################
        ppt_list = self.preprocess()
        ppt_list = [item for item in ppt_list if item.split(
            '.')[-1].lower() in ("ppt", "pptx")]
        try:
            for file_path in ppt_list:

                ext = os.path.splitext(file_path)[-1][1:]

                txt_file = (os.path.splitext(file_path)[0]) + ".txt"

                file_path = os.path.abspath(file_path)

                if ext == 'ppt':
                    file_path = self.convert_ppt_to_pptx(
                        file_path)  # First convert the ppt to pptx

                presentation = Presentation(file_path)

                logger.info(F"Slide count ={DMAGENTA} {
                            len(presentation.slides)}{RESET}")

                _slide_count_ = 0

                with Progress() as progress:

                    task = progress.add_task(
                        "[magenta]Preparing..", total=len(presentation.slides))

                    for slide in presentation.slides:

                        _slide_count_ += 1
                        # progress.console.print(F"Slide {_slide_count_}/{len(presentation.slides)}", end='\n')

                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame

                                for paragraph in text_frame.paragraphs:
                                    # Create a paragraph in the Word document if it contains text
                                    # Ensure text exists
                                    if any(run.text.strip() for run in paragraph.runs):

                                        for run in paragraph.runs:
                                            text = run.text.strip()
                                            if text and text != ' ':

                                                with open(txt_file, 'a') as fl:
                                                    fl.write(text)
                                                    # return txt_file

                        progress.update(task, advance=1)

                    if dest == "text":
                        with open(txt_file, "r") as fl:
                            text_buffer = fl.read()
                        print(text_buffer)
                        return text_buffer

                logger.info(f"{MAGENTA}New file is {CYAN}{txt_file}{RESET}")
                logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
        except Exception as e:
            logger.error(f"\n❌Oops! {RED}{e}{RESET}")

    @staticmethod
    def convert_ppt_to_pptx(obj: os.PathLike):
        import platform
        try:
            if obj.endswith('ppt'):
                if platform.system() in ('Linux', 'MacOS') or os.name == 'posix':
                    subprocess.run(
                        ['soffice', '--headless', '--convert-to', 'pptx', obj])
                    return os.path.splitext(obj)[0] + ".pptx"
                elif platform.system() in ('Windows') or os.name == 'nt':
                    import win32com.client
                    powerpoint = win32com.client.Dispatch(
                        "PowerPoint.Application")
                    powerpoint.Visible = 1
                    ppt = powerpoint.Presentations.Open(obj)
                    pptx_file = os.path.splitext(obj)[0] + ".pptx"
                    ppt.SaveAs(pptx_file, 24)  # 24 is the format for pptx
                    ppt.Close()
                    powerpoint.Quit()
                    return pptx_file
            else:
                print(F"{RED}Unable to identify the system{RESET}")
        except KeyboardInterrupt:
            print("\nQuit!")
            sys.exit(1)
        except Exception as e:
            logger.error(f"{RED}{e}{RESET}")

    def ppt_to_word(self):
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        from docx.shared import Pt
        from docx.shared import RGBColor as docxRGBColor
        from pptx.dml.color import RGBColor as pptxRGBColor

        ###############################################################################
        """Convert ppt file to word document
    ->Preserves bold formatting
    """
        ###############################################################################
        ppt_list = self.preprocess()
        ppt_list = [item for item in ppt_list if item.split(
            '.')[-1].lower() in ("ppt", "pptx")]
        for file_path in ppt_list:
            ext = os.path.splitext(file_path)[-1][1:]
            word_file = (os.path.splitext(file_path)[
                         0] + ".docx") if ext in list(_ext_ppt_) else None
            try:
                logger.info(f"{DYELLOW}Create Doument Tablet{RESET}")
                file_path = os.path.abspath(file_path)
                if ext == 'ppt':
                    file_path = self.convert_ppt_to_pptx(
                        file_path)  # First convert the ppt to pptx
                presentation = Presentation(file_path)
                document = Document()
                logger.info(
                    F"Slide count ={DMAGENTA} {len(presentation.slides)}{RESET}")
                _slide_count_ = 0
                with Progress() as progress:
                    task = progress.add_task(
                        "[magenta]Preparing..", total=len(presentation.slides))
                    for slide in presentation.slides:
                        _slide_count_ += 1
                        # progress.console.print(F"Slide {_slide_count_}/{len(presentation.slides)}", end='\n')
                        slide_text = ""
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame

                                for paragraph in text_frame.paragraphs:
                                    # Create a paragraph in the Word document if it contains text
                                    # Ensure text exists
                                    if any(run.text.strip() for run in paragraph.runs):
                                        # print("Has text")
                                        new_paragraph = document.add_paragraph()

                                        # Set general paragraph properties
                                        new_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY  # Justify text
                                        new_paragraph.space_after = Pt(6)
                                        new_paragraph.space_before = Pt(6)
                                        new_paragraph.line_spacing = 1.15

                                        for run in paragraph.runs:
                                            if run.text.strip():
                                                slide_text += run.text  # Only add non-empty text runs
                                                # print(run.text.strip(), end='\n')
                                                new_run = new_paragraph.add_run(
                                                    run.text)

                                                # Preserve bold, italic, underline, font name, and size
                                                new_run.bold = run.font.bold
                                                new_run.italic = run.font.italic
                                                new_run.underline = run.font.underline
                                                new_run.font.name = run.font.name
                                                new_run.font.size = run.font.size

                                                # Preserve font color
                                                try:
                                                    if run.font.color and run.font.color.rgb:
                                                        pptx_color = run.font.color.rgb
                                                        # If the color is white (255, 255, 255), change it to black (0, 0, 0)
                                                        if pptx_color == pptxRGBColor(255, 255, 255):
                                                            new_run.font.color.rgb = docxRGBColor(
                                                                0, 0, 0)  # Black
                                                        else:
                                                            # Assign color properly to the Word run
                                                            new_run.font.color.rgb = docxRGBColor(
                                                                pptx_color[0], pptx_color[1], pptx_color[2])
                                                except AttributeError:
                                                    pass

                        progress.update(task, advance=1)
                document.save(word_file)
                logger.info(f"{MAGENTA}New file is {CYAN}{word_file}{RESET}")
                logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
                return word_file
            except Exception as e:
                logger.error(
                    f"\n❌Oops! {RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(
                        f"\n❌Oops! {RED}{e}{RESET}")

    def text_to_word(self):

        ###############################################################################
        """Convert text file to word
    ->Read the text file
    ->Filter out non-XML characters
    ->Create a new Word document
    ->Add the filtered text content to the document"""
        ###############################################################################
        flist = self.preprocess()
        flist = [item for item in flist if item.lower().endswith("txt")]
        for file_path in flist:
            if file_path.lower().endswith("txt"):
                word_file = file_path[:-3] + "docx"

            try:
                # Read the text file
                logger.info(f"{DCYAN}Open and read the text file{RESET}")
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    text_content = file.read()

                # Filter out non-XML characters
                filtered_content = re.sub(
                    r'[^\x09\x0A\x0D\x20-\uD7FF\uE000-\uFFFD]+', '', text_content)

                # Create a new Word document
                logger.info(f"{DYELLOW}Create Doument Tablet{RESET}")
                doc = Document()
                # Add the filtered text content to the document
                doc.add_paragraph(filtered_content)

                # Save the document as a Word file
                doc.save(word_file)
                logger.info(f"{MAGENTA}New file is {DCYAN}{word_file}{RESET}")
                logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            except FileExistsError as e:
                logger.error(f"{str(e)}📁")
            except Exception as e:
                logger.error(
                    f"\n❌Oops something went awry {RED}{e}{RESET}")
                with open("conversion.log", "a") as log_file:
                    log_file.write(
                        f"\n❌Oops something went astray{RED}{e}{RESET}")

    def convert_xls_to_word(self):
        ###############################################################################
        """Convert xlsx file(s) to word file(s)
    ->Read the XLS file using pandas
    ->Create a new Word document
    ->Iterate over the rows of the dataframe and add them to the Word document"""
        ###############################################################################
        xls_list = self.preprocess()

        xls_list = [item for item in xls_list if item.split(
            '.')[-1].lower() in ("xls", "xlsx")]

        print(F"{DGREEN}Initializing conversion sequence{RESET}")

        for xls_file in xls_list:
            ext = os.path.splitext(xls_file)[-1][1:]
            word_file = (os.path.splitext(xls_file)[
                         0] + ".docx") if ext in list(_ext_xls) else None
            try:
                '''Read the XLS file using pandas'''

                df = pd.read_excel(xls_file)

                '''Create a new Word document'''
                doc = Document()

                '''Iterate over the rows of the dataframe and add them to the
                Word document'''
                logger.info(f"{ICYAN}Converting {xls_file}..{RESET}")
                # time.sleep(2)
                total_rows = df.shape[0]
                for _, row in df.iterrows():
                    current_row = _ + 1
                    percentage = (current_row / total_rows)*100
                    for value in row:
                        doc.add_paragraph(str(value))
                    print(f"Row {DYELLOW}{current_row}/{total_rows} \
{DBLUE}{percentage:.1f}%{RESET}", end="\r")
                    # print(f"\033[1;36m{row}{RESET}")

                # Save the Word document
                doc.save(word_file)
                print(F"{DGREEN}Conversion successful!{RESET}", end="\n")
            except KeyboardInterrupt:
                print("\nQuit⌨️")
                sys.exit(1)
            except Exception as e:
                print(f"{RED}Oops Conversion failed:❕{RESET}", str(e))

    def convert_xls_to_text(self):

        ###############################################################################
        '''Convert xlsx/xls file/files to text file format
    ->Read the XLS file using pandas
    ->Convert the dataframe to plain text
    ->Write the plain text to the output file'''
        ###############################################################################
        xls_list = self.preprocess()

        xls_list = [
            item for item in xls_list if any(item.lower().endswith(ext)
                                             for ext in _ext_xls)]
        print(F"{DGREEN}Initializing conversion sequence{RESET}")
        for xls_file in tqdm(xls_list):
            ext = os.path.splitext(xls_file)[-1][1:]
            txt_file = (os.path.splitext(xls_file)[
                        0] + ".txt") if ext in list(_ext_xls) else None
            try:
                # Read the XLS file using pandas
                logger.info(f"Converting {xls_file}..")
                df = pd.read_excel(xls_file)

                # Convert the dataframe to plain text
                text = df.to_string(index=False)
                chars = len(text)
                words = len(text.split())
                lines = len(text.splitlines())

                print(
                    f"Preparing to write: {DYELLOW}{chars} \033[1;30m \
characters{DYELLOW} {words}\033[1;30m words {DYELLOW}{lines}\033[1;30m \
lines {RESET}", end="\n")
                # Write the plain text to the output file
                with open(txt_file, 'w') as file:
                    file.write(text)

                print(F"{DGREEN}Conversion successful!{RESET}", end="\n")
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                print("Oops Conversion failed:", str(e))

    def convert_xlsx_to_csv(self):
        ###############################################################################
        '''Convert xlsx/xls file to csv(comma seperated values) format
        ->Load the Excel file
        ->Save the DataFrame to CSV'''
        ###############################################################################
        xls_list = self.preprocess()

        xls_list = [
            item for item in xls_list if item.split('.')[-1].lower() in ("xls", "xlsx")]
        for xls_file in tqdm(xls_list):
            ext = os.path.splitext(xls_file)[-1][1:]
            csv_file = (os.path.splitext(xls_file)[
                        0] + ".csv") if ext in list(_ext_xls) else None
            try:
                '''Load the Excel file'''
                print(F"{DGREEN}Initializing conversion sequence{RESET}")
                df = pd.read_excel(xls_file)
                logger.info(f"Converting {xls_file}..")
                total_rows = df.shape[0]
                print(f"Writing {DYELLOW}{total_rows} rows {RESET}", end="\n")
                for i in range(101):
                    print(f"Progress: {i}%", end="\r")
                '''Save the DataFrame to CSV'''
                df.to_csv(csv_file, index=False)
                print(F"{DMAGENTA} Conversion successful{RESET}")

            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                print(e)

    def convert_csv_to_xlsx(self):
        csv_list = self.preprocess()
        csv_list = [
            item for item in csv_list if item.split('.')[-1].lower() in ("csv")]

        with Progress() as progress:
            task = progress.add_task("[cyan]Coverting", total=len(csv_list))
            for file in csv_list:
                file_name = file[:-3] + 'xlsx'
                df = pd.read_csv(file)
                # excel engines ('openpyxl' or 'xlsxwriter')
                df.to_excel(file_name, engine='openpyxl', index=False)

                # Load the workbook and the sheet
                workbook = load_workbook(file_name)
                sheet = workbook.active

                # print("Adjust Columns")
                for column in sheet.columns:
                    max_length = 0
                    column_letter = column[0].column_letter

                    max_length = max(len(str(cell.value)) for cell in column)
                    '''for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(cell.value)

                        except Exception:
                            pass
                    '''

                    adjusted_width = (max_length + 2)
                    sheet.column_dimensions[column_letter].width = adjusted_width

                # Save the workbook
                workbook.save(file_name)
                progress.update(task, advance=1)

    def convert_xlsx_to_database(self):
        ###############################################################################
        """Convert xlsx file(s) to sqlite
        ->Read the Excel file into a pandas DataFrame
        ->Create a connection to the SQLite database
        ->Insert the DataFrame into a new table in the database
        ->Close the database connection"""
        ###############################################################################
        xlsx_list = self.preprocess()
        xlsx_list = [
            item for item in xlsx_list if item.split('.')[-1].lower() in ("xls", "xlsx")]
        for xlsx_file in tqdm(xlsx_list):
            sqlfile = (os.path.splitext(xlsx_file)[
                       0] + ".sql") if (xlsx_file.split('.')[0]) in ("xls", "xlsx") else None
            try:
                db_file = input(
                    F"{DBLUE}Please enter desired sql filename: {RESET}")
                table_name = input(
                    "Please enter desired table name: ")
                # res = ["db_file", "table_name"]
                if any(db_file) == "":
                    db_file = sqlfile
                    table_name = sqlfile[:-4]
                if not db_file.endswith(".sql"):
                    db_file = sqlfile
                column = 0
                for i in range(20):
                    column += 0
                # Read the Excel file into a pandas DataFrame
                print(f"Reading {xlsx_file}...")
                df = pd.read_excel(xlsx_file)
                print(f"{DGREEN}Initializing conversion sequence{RESET}")
                print(f"{DGREEN} Connected to sqlite3 database::{RESET}")
                # Create a connection to the SQLite database
                conn = sqlite3.connect(db_file)
                print(F"{DYELLOW} Creating database table::{RESET}")
                # Insert the DataFrame into a new table in the database
                df.to_sql(table_name, column, conn,
                          if_exists='replace', index=False)
                print(
                    f"Operation successful{RESET} file saved as \033[32{db_file}{RESET}")
                # Close the database connection
                conn.close()
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)
            except Exception as e:
                logger.error(f"{e}")

    def doc2image(self, outf="png"):
        ###############################################################################
        """Create image objects from given files"""
        ###############################################################################
        outf = "png" if outf not in ('png', 'jpg') else outf
        path_list = self.preprocess()
        file_list = [
            item for item in path_list if item.split('.')[-1].lower() in ("pdf", "doc", "docx")]
        imgs = []
        for file in file_list:
            if file.lower().endswith("pdf"):
                # Convert the PDF to a list of PIL image objects
                print(f"{DBLUE}Generate image objects ..{RESET}")
                images = convert_from_path(file)

                # Save each image to a file
                fname = file[:-4]
                print(f"{YELLOW}Target images{BLUE} {len(images)}{RESET}")

                with Progress() as progress:
                    task = progress.add_task(
                        "[magenta]Generating images ", total=len(images))
                    for i, image in enumerate(images):
                        # print(f"{DBLUE}{i}{RESET}", end="\r")
                        yd = f"{fname}_{i+1}.{outf}"
                        image.save(yd)
                        imgs.append(yd)
                        progress.update(task, advance=1)
                # print(f"\n{GREEN}Ok{RESET}")

        return imgs


class Scanner:
    """Implementation of scanning to extract data from pdf files and images
    input_file -> file to be scanned pdf,image
    Args:
        input_file->file to be scanned
        no_strip-> Preserves text formating once set to True, default: False
    Returns:
        None"""

    def __init__(self, input_file, no_strip: bool = False):
        self.input_file = input_file
        self.no_strip = no_strip

    def preprocess(self):
        files_to_process = []

        if os.path.isfile(self.input_file):
            files_to_process.append(self.input_file)
        elif os.path.isdir(self.input_file):
            for file in os.listdir(self.input_file):
                file_path = os.path.join(self.input_file, file)
                if os.path.isfile(file_path):
                    files_to_process.append(file_path)

        return files_to_process

    def scanPDF(self, obj=None):
        """Obj - object for scanning where the object is not a list"""
        pdf_list = self.preprocess()
        pdf_list = [item for item in pdf_list if item.lower().endswith("pdf")]
        if obj:
            pdf_list = [obj]

        for pdf in pdf_list:
            out_f = pdf[:-3] + 'txt'
            print(f"{YELLOW}Read pdf ..{RESET}")

            with open(pdf, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ''

                pg = 0
                for page_num in range(len(reader.pages)):
                    pg += 1

                    print(f"{DYELLOW}Progress:{RESET}", end="")
                    print(f"{CYAN}{pg}/{len(reader.pages)}{RESET}", end="\r")
                    page = reader.pages[page_num]
                    text += page.extract_text()

            print(f"\n{text}")
            print(F"\n{YELLOW}Write text to {GREEN}{out_f}{RESET}")
            with open(out_f, 'w') as f:
                f.write(text)

            print(F"\n{DGREEN}Ok{RESET}")

    def scanAsImgs(self):
        file = self.input_file
        mc = MakeConversion(file)
        img_objs = mc.doc2image()
        # print(img_objs)
        from .OCRTextExtractor import ExtractText
        text = ''

        with Progress() as progress:
            task = progress.add_task(
                "[magenta]Extracting text", totsl=len(img_objs))
            for i in img_objs:
                extract = ExtractText(i, self.no_strip)
                _text = extract.OCR()

                if _text is not None:
                    text += _text
                    with open(f"{self.input_file[:-4]}_filemac.txt", 'a') as _writer:
                        _writer.write(text)
                progress.update(task, advance=1)

        def _cleaner_():
            print(f"{FMAGENTA}Clean")
            for obj in img_objs:
                if os.path.exists(obj):
                    print(obj, end='\r')
                    os.remove(obj)
                txt_file = f"{obj[:-4]}.txt"
                if os.path.exists(txt_file):
                    print(f"{FCYAN}{txt_file}{RESET}", end='\r')
                    os.remove(txt_file)

        # Do clean up
        _cleaner_()
        from .overwrite import clear_screen
        clear_screen()
        print(f"{DGREEN}{IGREEN}Full Text{RESET}")
        print(text)
        print(
            f"{BWHITE}Text File ={IGREEN}{self.input_file[:-4]}_filemac.txt{RESET}")
        print(f"{GREEN}Ok✅{RESET}")
        return text

    def scanAsLongImg(self) -> bool:
        try:
            """Convert the pdf to long image for scanning - text extraction"""
            pdf_list = self.preprocess()
            pdf_list = [
                item for item in pdf_list if item.lower().endswith("pdf")]
            from .longImg import LImage
            from .OCRTextExtractor import ExtractText
            for file in pdf_list:
                LI = LImage(file)
                fl = LI.preprocess()

                # fpath = file.split('.')[0] + '.png'
                tx = ExtractText(fl, self.no_strip)
                text = tx.OCR()
                if text is not None:
                    # print(text)
                    print(f"{GREEN}Ok{RESET}")
            return True
        except Exception as e:
            print(e)


class FileSynthesis:
    '''Definition of audiofying class'''

    def __init__(self, obj: Union[os.PathLike, str, List[Union[os.PathLike, str]]], resume: bool = True):
        self.obj = obj
        self.resume = resume

    @staticmethod
    def join_audios(files, output_file):
        masterfile = output_file + "_master.mp3"
        print(
            f"{DBLUE}Create a master file {DMAGENTA}{masterfile}{RESET}", end='\r')
        # Create a list to store files
        ogg_files = []
        # loop through the directory while adding the ogg files to the list
        for filename in files:
            print(f"Join {DBLUE}{len(files)}{RESET} files")
            # if filename.endswith('.ogg'):
            # ogg_file = os.path.join(path, filename)
            ogg_files.append(AudioSegment.from_file(filename))

        # Concatenate the ogg files
        combined_ogg = ogg_files[0]
        for i in range(1, len(files)):
            combined_ogg += ogg_files[i]

        # Export the combined ogg to new mp3 file or ogg file
        combined_ogg.export(output_file + "_master.ogg", format='ogg')
        print(F"{DGREEN}Master file:Ok                                                                             {RESET}")

    def Synthesise(self, text: str, output_file: str, CHUNK_SIZE: int = 1_000, _tmp_folder_: str = 'tmp_dir', thread_name: str = None, max_retries: int = 30) -> None:
        """Converts given text to speech using Google Text-to-Speech API."""
        # from rich.progress import (BarColumn, Progress, SpinnerColumn,TextColumn)

        config = ConfigManager()
        # Define directories and other useful variables for genrating output_file and checkpoint_file
        out_dir = os.path.split(output_file)[0]

        thread_name = f"thread_{os.path.split(output_file.split('.')[0])[-1]}"
        _file_ = os.path.split(output_file)[1]

        _tmp_folder_ = os.path.join(out_dir, _tmp_folder_)

        # Remove temporary dir if it exists, rare-cases since file names are mostly unique
        if os.path.exists(_tmp_folder_) and self.resume is False:
            # query = input(f"{DBLUE}Remove the {os.path.join(out_dir, _tmp_folder_)} directory (y/n)?{RESET} ").lower() in ('y', 'yes')
            shutil.rmtree(_tmp_folder_)

        # Create temporary folder to house chunks
        if not os.path.exists(_tmp_folder_):
            logger.info(
                f"{DYELLOW}Create temporary directory = {DBLUE}{_tmp_folder_}{RESET}")
            os.mkdir(_tmp_folder_)

        _full_output_path_ = os.path.join(_tmp_folder_, _file_)

        ''' Create a checkpoint_file for conversion resumption if need be
        _check_file = os.path.splitext(_file_)[0] + '.ch'

        _check_dir_path_ = os.path.join(out_dir, _tmp_folder_, _check_file)

        checkpoint_file = _check_dir_path_

        start_chunk = 0  # Initialize start chank

        Check if there is any checkpoint record in chekpoint file for resumption
        if os.path.exists(checkpoint_file):
            logger.info(f"{DYELLOW}Found a Checkpoint file{RESET}")
            with open(checkpoint_file, 'r') as f:
                start_chunk = int(f.read())
            logger.info(
                f"{DYELLOW}Resuming from chunk{DBLUE} {start_chunk}{RESET}")'''

        # Read reume chunk from the configuration file
        start_chunk = int(config.read_config_file(thread_name)) * 1_000
        start_chunk = 0 if start_chunk is None else start_chunk

        ''' If chunk is not 0 multiply the chunk by the highest decimal value of the chunk size
        else set it to 0 meaning file is being operated on for the first time
        '''
        resume_chunk_pos = start_chunk * 1_000 if start_chunk != 0 else start_chunk

        try:

            print(f"{DYELLOW}Start thread:: {thread_name}{RESET}")

            total_chunks = math.ceil(len(text) / CHUNK_SIZE)

            counter = math.ceil(resume_chunk_pos /
                                CHUNK_SIZE) if resume_chunk_pos != 0 else 0

            attempt = 0

            while attempt <= max_retries:
                try:
                    # Initialize progress bar for the overall process

                    for i in range(resume_chunk_pos, len(text), CHUNK_SIZE):

                        print(
                            f"Processing: chunk {MAGENTA}{counter}/{total_chunks} {DCYAN}{counter/total_chunks*100:.2f}%{RESET}\n", end="\r")
                        chunk = text[i:i + CHUNK_SIZE]
                        # print(chunk)
                        if os.path.exists(f"{_full_output_path_}_{counter}.ogg"):

                            if counter == start_chunk:
                                print(
                                    f"{CYAN}Chunk vs file confict: {BLUE}Resolving{RESET}")
                                os.remove(f"{_full_output_path_}_{
                                    counter}.ogg")
                                output_filename = f"{
                                    _full_output_path_}_{counter}.ogg"

                            # Remove empty file
                            elif os.path.getsize(f"{_full_output_path_}_{counter}.ogg") != 0:
                                os.remove(f"{_full_output_path_}_{
                                    counter}.ogg")
                                output_filename = f"{
                                    _full_output_path_}_{counter}.ogg"

                            else:
                                output_filename = f"{
                                    _full_output_path_}_{counter+1}.ogg"

                        else:
                            output_filename = f"{
                                _full_output_path_}_{counter}.ogg"

                        tts = gTTS(text=chunk, lang='en', slow=False)

                        tts.save(output_filename)

                        # Update current_chunk in the configuration
                        config.update_config_entry(
                            thread_name, current_chunk=counter)

                        counter += 1

                except FileNotFoundError as e:
                    logger.error(f"{RED}{e}{RESET}")

                except requests.exceptions.ConnectionError:  # Handle connectivity/network error
                    logger.error(f"{RED}ConnectionError{RESET}")

                    # Exponential backoff for retries
                    for _sec_ in range(2 ** attempt, 0, -1):
                        print(
                            # Increament the attempts
                            f"{BWHITE}Resume in {DBLUE}{_sec_}{RESET}", end='\r')

                    attempt += 1

                    # Read chunk from configuration
                    resume_chunk_pos = int(
                        config.read_config_file(thread_name)) * 1_000

                except requests.exceptions.HTTPError as e:  # Exponential backoff for retries
                    logger.error(f"HTTP error: {e.status_code} - {e.reason}")
                    for _sec_ in range(2 ** attempt, 0, -1):
                        print(
                            f"{BWHITE}Resume in {DBLUE}{_sec_}{RESET}", end='\r')

                    attempt += 1

                    resume_chunk_pos = int(
                        config.read_config_file(thread_name)) * 1_000

                except requests.exceptions.RequestException as e:
                    logger.error(f"{RED}{e}{RESET}")

                    for _sec_ in range(2 ** attempt, 0, -1):
                        print(
                            f"{BWHITE}Resume in {DBLUE}{_sec_}{RESET}", end='\r')

                    attempt += 1

                    resume_chunk_pos = int(
                        config.read_config_file(thread_name)) * 1_000

                except (ConnectionError, ConnectionAbortedError, ConnectionRefusedError, ConnectionResetError):
                    logger.error(
                        f'{RED}Connection at attempt{RESET}')

                    for _sec_ in range(2 ** attempt, 0, -1):
                        print(
                            f"{BWHITE}Resume in {DBLUE}{_sec_}{RESET}", end='\r')

                        attempt += 1

                    resume_chunk_pos = int(
                        config.read_config_file(thread_name)) * 1_000

                except MarkupError as e:
                    logger.error(F"{RED}{e}{RESET}")
                except Exception as e:  # Handle all other types of exceptions
                    logger.error(
                        f'{DMAGENTA}{attempt+1}/{max_retries}:{RED}{e}{RESET}')

                    for _sec_ in range(2 ** attempt, 0, -1):
                        pass

                    attempt += 1

                    resume_chunk_pos = int(
                        config.read_config_file(thread_name)) * 1_000

                else:
                    print(
                        f"{FMAGENTA}Conversion success✅. \n  {FCYAN}INFO\t Create masterfile{RESET}")

                    if len(os.listdir(_tmp_folder_)) > 2:  # Combine generated gTTS objects

                        from .JoinAudios import JoinAudios

                        joiner = JoinAudios(
                            _tmp_folder_, masterfile=output_file)
                        joiner.worker()
                        # Remove temporary files
                        shutil.rmtree(_tmp_folder_)

                    break  # Exit the retry loop if successfull

            else:
                print(
                    f"{RED}Maximum retries reached. Unable to complete the operation after {DMAGENTA} {max_retries} attempts.{RESET}")
                sys.exit(2)

        finally:
            pass

    @staticmethod
    def pdf_to_text(pdf_path):
        logger.info(
            F'{GREEN} Initializing pdf to text conversion{RESET}')
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ''
                _pg_ = 0
                print(F"{YELLOW}Convert pages..{RESET}")
                for page_num in range(len(pdf_reader.pages)):
                    _pg_ += 1
                    logger.info(
                        f"Page {DBLUE}{_pg_}{RESET}/{len(pdf_reader.pages)}")
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()
                print(F"{DGREEN}Ok{RESET}\n")
                return text
        except Exception as e:
            logger.error(
                f"{DRED}Failed to extract text from '{YELLOW}{pdf_path}'{RESET}:\n {e}")

    @staticmethod
    def text_file(input_file):
        try:
            with open(input_file, 'r', errors='ignore') as file:
                text = file.read().replace('\n', ' ')
            return text
        except FileNotFoundError:
            logger.error("File '{}' was not found.📁".format(input_file))
        except Exception as e:
            logger.error(
                F"{DRED}{str(e)}{RESET}")

    @staticmethod
    def docx_to_text(docx_path):
        try:
            logger.info(f"{BLUE} Converting {docx_path} to text{RESET}")
            doc = Document(docx_path)
            paragraphs = [paragraph.text for paragraph in doc.paragraphs]
            return '\n'.join(paragraphs)
        except FileNotFoundError:
            logger.error(f"File '{docx_path}' was not found.📁")
        except Exception as e:
            logger.error(
                F"{DRED}Error converting {docx_path} to text: {e}\
{RESET}")

    class THAudio:
        def __init__(self, instance):
            self.instance = instance
            self.lock = Lock()
            self.config = ConfigManager()

        def audiofy(self, num_threads=3):
            ls = ("pdf", "docx", "doc", "txt", "ppt", "pptx")

            def create_thread(item, thread_name):
                # Create a unique temp dir for each file
                temp_dir = f"tmp_dir_{os.path.split(item.split('.')[0])[-1]}"

                # Ensure proper locking when adding config entry
                with self.lock:
                    # Record config entry for each item
                    self.config.add_config_entry(
                        thread_name, f"{item.split('.')[0]}", temp_dir, 0)

                # Create and return the thread
                return Thread(target=self.worker, args=(item, temp_dir, thread_name), name=thread_name)

            threads = []
            processed_items = 0

            # Process a list of files
            def process_batch():
                for thread in threads:
                    thread.start()
                for thread in threads:
                    thread.join()
                threads.clear()  # Clear thread list after batch is done

            # Handle files provided as a list
            if isinstance(self.instance.obj, list):
                for item in self.instance.obj:
                    item = os.path.abspath(item)
                    if os.path.isfile(item) and item.endswith(ls):
                        thread_name = f"thread_{os.path.split(item.split('.')[0])[-1]}"
                        thread = create_thread(item, thread_name)
                        threads.append(thread)
                        processed_items += 1

                        # Process threads in batches of 'num_threads'
                        if processed_items % num_threads == 0:
                            process_batch()

                # Process remaining threads in case the list isn't a perfect multiple of num_threads
                if threads:
                    process_batch()

            # Handle a single file
            elif os.path.isfile(self.instance.obj):
                item = os.path.abspath(self.instance.obj)
                if item.endswith(ls):
                    thread_name = f"thread_{os.path.split(item.split('.')[0])[-1]}"
                    thread = create_thread(item, thread_name)
                    threads.append(thread)
                    process_batch()  # Process immediately for single file

            # Handle a directory of files
            elif os.path.isdir(self.instance.obj):
                for item in os.listdir(self.instance.obj):
                    item = os.path.abspath(item)
                    if os.path.isfile(item) and item.endswith(ls):
                        thread_name = f"thread_{os.path.split(item.split('.')[0])[-1]}"
                        thread = create_thread(item, thread_name)
                        threads.append(thread)
                        processed_items += 1

                        # Process threads in batches
                        if processed_items % num_threads == 0:
                            process_batch()

                # Process remaining threads
                if threads:
                    process_batch()

        def worker(self, input_file, _temp_dir_, thread_name):
            output_file = os.path.split(input_file)[-1].split('.')[0] + '.ogg'
            print(f"Thread {thread_name} processing file: {input_file}")

            try:
                # Extract text based on file type
                if input_file.endswith('.pdf'):
                    text = FileSynthesis.pdf_to_text(input_file)
                elif input_file.lower().endswith(tuple(_ext_word)):
                    text = FileSynthesis.docx_to_text(input_file)
                elif input_file.endswith('.txt'):
                    text = FileSynthesis.text_file(input_file)
                elif input_file.split('.')[-1] in ("ppt", "pptx"):
                    conv = MakeConversion(input_file)
                    word = conv.ppt_to_word()
                    conv = MakeConversion(word)
                    text = FileSynthesis.text_file(conv.word_to_txt())
                else:
                    raise ValueError(
                        'Unsupported file format. Please provide a PDF, txt, or Word document.')

                # Synthesize audio using the extracted text
                self.instance.Synthesise(
                    text, output_file, _tmp_folder_=_temp_dir_, thread_name=thread_name)
                print(
                    f"Thread {thread_name} completed processing {input_file}")

            except Exception as e:
                print(f"Error in thread {thread_name}: {e}")
            except KeyboardInterrupt:
                print(f"Thread {thread_name} interrupted.")
                sys.exit(1)


class ConfigManager:
    def __init__(self, config_path='filemac_config.json'):
        self.config_path = config_path

    def create_config_file(self, config_data):
        """
        Create or overwrite a configuration file to record thread names, associated file names, and current chunks.

        Args:
            config_data(list): A list of dictionaries containing thread name, associated file name, temp dir, and current chunk.
        """
        try:
            # Ensure the output directory exists
            output_dir = os.path.dirname(self.config_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # Write the configuration data to a JSON file
            with open(self.config_path, 'w') as config_file:
                json.dump(config_data, config_file, indent=4)

            print(
                f"Configuration file '{self.config_path}' created successfully.")
        except Exception as e:
            print(f"Error creating configuration file: {e}")

    def read_config_file(self, thread=None):
        """
        Read the configuration file and return the data or a specific thread's current chunk.

        Args:
            thread (str): The thread name to search for in the config. If None, returns the full config.

        Returns:
            dict/list: Returns a specific entry for the thread or the full configuration data.
            None: If the file doesn't exist or thread is not found.
        """
        try:
            if not os.path.exists(self.config_path):
                print(f"Configuration file '{self.config_path}' not found.")
                return None

            with open(self.config_path, 'r') as config_file:
                config = json.load(config_file)

            if thread is None:
                return config  # Return entire configuration

            # Search for specific thread's current chunk
            for entry in config:
                if entry['thread_name'] == thread:
                    return entry.get('current_chunk', None)

            print(f"Entry for thread '{thread}' not found.")
            return None

        except Exception as e:
            print(f"Error reading configuration file: {e}")
            return None

    def add_config_entry(self, thread_name, associated_file, tmp_dir, current_chunk):
        """
        Add a new entry to the configuration file.

        Args:
            thread_name (str): The name of the thread to be added.
            associated_file (str): The associated file name for the thread.
            tmp_dir (str): Temporary directory for the thread.
            current_chunk (int): The current chunk number for the thread.
        """
        try:
            # Read existing config data or create a new list if the file doesn't exist
            config_data = self.read_config_file() or []

            # Check if the thread already exists in the configuration
            for entry in config_data:
                if entry['thread_name'] == thread_name:
                    print(
                        f"Thread '{thread_name}' already exists. Use 'update_config_entry' to update it.")
                    return

            # Add the new entry
            config_data.append({
                "thread_name": thread_name,
                "associated_file": associated_file,
                "tmp_dir": tmp_dir,
                "current_chunk": current_chunk
            })

            # Save the updated configuration
            self.create_config_file(config_data)

        except Exception as e:
            print(f"Error adding config entry: {e}")

    def update_config_entry(self, thread_name, associated_file=None, tmp_dir=None, current_chunk=None):
        """
        Update an existing entry in the configuration file.

        Args:
            thread_name (str): The name of the thread to update.
            associated_file (str, optional): The updated associated file name. Defaults to None.
            tmp_dir (str, optional): The updated temporary directory. Defaults to None.
            current_chunk (int, optional): The updated current chunk number. Defaults to None.
        """
        try:
            # Read existing config data
            config_data = self.read_config_file() or []

            # Find the entry to update
            for entry in config_data:
                if entry['thread_name'] == thread_name:
                    if associated_file:
                        entry['associated_file'] = associated_file
                    if tmp_dir:
                        entry['tmp_dir'] = tmp_dir
                    if current_chunk is not None:
                        entry['current_chunk'] = current_chunk

                    # Save the updated configuration
                    self.create_config_file(config_data)
                    print(f"Thread '{thread_name}' updated successfully.")
                    return True

            print(f"Thread '{thread_name}' not found in the configuration.")

        except Exception as e:
            print(f"Error updating config entry: {e}")


###############################################################################
# Convert video file to from one format to another'''
###############################################################################


class VideoConverter:

    def __init__(self, input_file, out_format=None):
        self.input_file = input_file
        self.out_format = out_format

    def preprocess(self):
        if self.out_format is None:
            return None
        files_to_process = []

        if os.path.isfile(self.input_file):
            files_to_process.append(self.input_file)
        elif os.path.isdir(self.input_file):
            if os.listdir(self.input_file) is None:
                print(f"{RED}Cannot work with empty folder{RESET}")
                sys.exit(1)
            for file in os.listdir(self.input_file):
                file_path = os.path.join(self.input_file, file)
                if os.path.isfile(file_path):
                    files_to_process.append(file_path)

        return files_to_process

    def ffmpeg_merger(self, obj: list = None):
        video_list = self.preprocess(), obj
        for input_video in video_list:
            output_file = [f"{_}_new_.{ext}" for _,
                           ext in [input_video.split('.', 1)]][0]
            # keep the original video quality by using -c:v copy, which avoids re-encoding.
            subprocess.run([
                "ffmpeg", "-i", input_video, "-i", "audio.mp3",
                "-c:v", "copy", "-c:a", "aac", "-strict", "experimental", "output_video.mp4"
            ])

    def pydub_merger(self, obj: list = None):
        video_list = self.preprocess() or obj
        for input_video in video_list:
            output_file = [f"{_}_new_.{ext}" for _,
                           ext in [input_video.split('.', 1)]][0]
            # Process or manipulate audio with Pydub (e.g., adjust volume)
            audio = AudioSegment.from_file("audio.mp3")
            audio = audio + 6  # Increase volume by 6 dB
            audio.export("processed_audio.mp3", format="mp3")

            # Merge processed audio with video using FFmpeg
            subprocess.run([
                "ffmpeg", "-i", input_video, "-i", "processed_audio.mp3",
                "-c:v", "copy", "-c:a", "aac", output_file
            ])

    def cv2_merger(self, obj: list = None):
        video_list = self.preprocess(), obj
        for input_video in video_list:
            # Read video and save frames (without audio)
            cap = cv2.VideoCapture(input_video)

            # Retrieve width and height from the video
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            fps = int(cap.get(cv2.CAP_PROP_FPS))

            # _, ext = input_video.split('.')[0]
            # output_file = f"{_}_new{ext}"
            output_file = [f"{_}_new_.{ext}" for _,
                           ext in [input_video.split('.', 1)]][0]
            # Define the VideoWriter with the video dimensions
            out = cv2.VideoWriter(output_file, cv2.VideoWriter_fourcc(
                *'mp4v'), fps, (width, height))

            # Read frames from the original video and write them to the output
            while cap.isOpened():
                ret, frame = cap.read()
                if not ret:
                    break
                out.write(frame)

            # Release resources
            cap.release()
            out.release()

            # Merge with audio using FFmpeg
            subprocess.run([
                "ffmpeg", "-i", "video_no_audio.mp4", "-i", "audio.mp3",
                "-c:v", "copy", "-c:a", "aac", "output_video.mp4"
            ])

    def CONVERT_VIDEO(self):
        try:
            input_list = self.preprocess()
            out_f = self.out_format.upper()
            input_list = [item for item in input_list if any(
                item.upper().endswith(ext) for ext in SUPPORTED_VIDEO_FORMATS)]
            print(F"{DYELLOW}Initializing conversion..{RESET}")

            for file in tqdm(input_list):
                if out_f.upper() in Video_codecs.keys():
                    _, ext = os.path.splitext(file)
                    output_filename = _ + '.' + out_f.lower()
                    # print(output_filename)
                elif out_f.upper() in SUPPORTED_VIDEO_FORMATS and not out_f.upper() in Video_codecs.keys():
                    print(
                        f"{RED}Unsupported output format --> Pending Implementation{RESET}")
                    sys.exit(1)
                else:
                    print(f"{RED}Unsupported output format{RESET}")
                    sys.exit(1)

                '''Load the video file'''
                print(f"{DBLUE}Load file{RESET}")
                video = VideoFileClip(file)
                '''Export the video to a different format'''
                print(f"{DMAGENTA}Converting file to {output_filename}{RESET}")
                video.write_videofile(
                    output_filename, codec=Video_codecs[out_f])
                '''Close the video file'''
                print(f"{DGREEN}Done{RESET}")
                video.close()
        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except Exception as e:
            print(e)


###############################################################################
# Convert Audio file to from one format to another'''
###############################################################################


class AudioConverter:

    def __init__(self, input_file, out_format):
        self.input_file = input_file
        self.out_format = out_format

    def preprocess(self):
        files_to_process = []

        if os.path.isfile(self.input_file):
            files_to_process.append(self.input_file)
        elif os.path.isdir(self.input_file):
            if os.listdir(self.input_file) is None:
                print(f"{RED}Cannot work with empty folder{RESET}")
                sys.exit(1)
            for file in os.listdir(self.input_file):
                file_path = os.path.join(self.input_file, file)
                if os.path.isfile(file_path):
                    files_to_process.append(file_path)

        return files_to_process

    def pydub_conv(self):
        try:
            input_list = self.preprocess()
            out_f = self.out_format
            input_list = [item for item in input_list if any(
                item.lower().endswith(ext) for ext in SUPPORTED_AUDIO_FORMATS)]
            print(F"{DYELLOW}Initializing conversion..{RESET}")

            def wav_redudancy():
                # Load the mp3 file using Pydub
                audio = pydub.AudioSegment.from_file(file, fmt)
                # Export the audio to a temporary file in wav format (ffmpeg can convert from wav to m4a)
                audio.export("temp.wav", format="wav")

            for file in tqdm(input_list):
                if out_f.lower() in SUPPORTED_AUDIO_FORMATS_DIRECT:
                    _, ext = os.path.splitext(file)
                    output_filename = _ + '.' + out_f
                    fmt = ext[1:]
                    # print(fmt, out_f)
                    audio = pydub.AudioSegment.from_file(file, fmt)
                    print(f"{DMAGENTA}Converting to {output_filename}{RESET}")
                    audio.export(output_filename, format=out_f)
                    # new_audio = pydub.AudioSegment.from_file('output_audio.')
                    print(f"{DGREEN}Done{RESET}")

                elif file[-3:].lower() == 'm4a' or out_f.lower() == "m4a":
                    _m4a_main_(file, out_f)

                elif out_f.lower() in SUPPORTED_AUDIO_FORMATS and not SUPPORTED_AUDIO_FORMATS_DIRECT:
                    print("Pending Implemantation For the format")

                else:
                    print(F"{RED}Unsupported output format{RESET}")
                    sys.exit(1)

        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except Exception as e:
            print(f"{RED}{e}{RED}")


###############################################################################
# Convert images file to from one format to another
###############################################################################


class ImageConverter:

    def __init__(self, input_file, out_format):
        self.input_file = input_file
        self.out_format = out_format

    def preprocess(self):
        try:
            files_to_process = []

            if os.path.isfile(self.input_file):
                files_to_process.append(self.input_file)
            elif os.path.isdir(self.input_file):
                if os.listdir(self.input_file) is None:
                    print("Cannot work with empty folder")
                    sys.exit(1)
                for file in os.listdir(self.input_file):
                    file_path = os.path.join(self.input_file, file)
                    if os.path.isfile(file_path):
                        files_to_process.append(file_path)

            return files_to_process
        except FileNotFoundError:
            print("File not found❕")
            sys.exit(1)

    def convert_image(self):
        try:
            input_list = self.preprocess()
            out_f = self.out_format.upper()
            input_list = [item for item in input_list if any(
                item.lower().endswith(ext) for ext in SUPPORTED_IMAGE_FORMATS.values())]

            for file in tqdm(input_list):
                if out_f.upper() in SUPPORTED_IMAGE_FORMATS:
                    _ = os.path.splitext(file)[0]
                    output_filename = _ + \
                        SUPPORTED_IMAGE_FORMATS[out_f].lower()
                else:
                    print("Unsupported output format")
                    sys.exit(1)
                '''Load the image using OpenCV: '''
                print(F"{DYELLOW}Reading input image..{RESET}")
                img = cv2.imread(file)
                '''Convert the OpenCV image to a PIL image: '''
                print(f"{DMAGENTA}Converting to PIL image{RESET}")
                pil_img = Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
                '''Save the PIL image to a different format: '''
                print(f"\033[1;36mSaving image as {output_filename}{RESET}")
                pil_img.save(output_filename, out_f)
                print(f"{DGREEN}Done ✅{RESET}")
                '''Load the image back into OpenCV: '''
                # print(f"{DMAGENTA}Load and display image{RESET}")
                # opencv_img = cv2.imread(output_filename)
                '''Display the images: '''
                # cv2.imshow('OpenCV Image', opencv_img)
                # opencv_img.show()
                '''Wait for the user to press a key and close the windows: '''
                # cv2.waitKey(0)
                # cv2.destroyAllWindows()
        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except AssertionError:
            print("Assertion failed.")
        except KeyError:
            print(
                f"{RED}ERROR:\tPending Implementation for{ICYAN} {out_f} {BWHITE}format{RESET}")
        except Exception as e:
            print(f"{RED}{e}{RESET}")
