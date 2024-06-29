from pptx import Presentation
from docx import Document
import sys


def pptx(ppt_file):
    try:
        print("Create Doument Tablet")
        presentation = Presentation(ppt_file)
        document = Document()
        _out_ = ppt_file.split('.')[0] + "docx"
        _slide_count_ = 0
        for slide in presentation.slides:
            _slide_count_ += 1
            print(f"INFO\t Slide {_slide_count_}/{len(presentation.slides)}", end='\r')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        new_paragraph = document.add_paragraph()
                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run(run.text)
                            # Preserve bold formatting
                            new_run.bold = run.font.bold
                            # Preserve italic formatting
                            new_run.italic = run.font.italic
                            # Preserve underline formatting
                            new_run.underline = run.font.underline
                            # Preserve font name
                            new_run.font.name = run.font.name
                            # Preserve font size
                            new_run.font.size = run.font.size
                            try:
                                # Preserve font color
                                new_run.font.color.rgb = run.font.color.rgb
                            except AttributeError:
                                # Ignore error and continue without
                                # setting the font color
                                pass
                    # Add a new paragraph after each slide
                    document.add_paragraph()
        document.save(_out_)
    except FileNotFoundError as e:
        print(e)

    except KeyboardInterrupt:
        print("\nQuit!")
        sys.exit(1)

    except Exception as e:
        print(e)
        pass


pptx(ppt_file="/home/skye/Software Engineering/Y1/SEM2/Data comm. networks/pdfs/10Data communication lecture 10.pptx")
